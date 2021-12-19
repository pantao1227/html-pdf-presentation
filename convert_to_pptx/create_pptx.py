import sys
from pptx import Presentation
from pptx.util import Inches

img_path = sys.argv[1:]

prs = Presentation('template16x9.pptx')
blank_slide_layout = prs.slide_layouts[0]

for i in img_path:
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(i, left = Inches(0), top = Inches(0), height = prs.slide_height)

prs.save('out.pptx')
